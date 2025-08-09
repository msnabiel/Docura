# pptx_ocr_extractor.py

import os
import re
import subprocess
import tempfile
from typing import Optional
from pptx import Presentation
import fitz  # PyMuPDF
import google.generativeai as genai
from concurrent.futures import ThreadPoolExecutor, as_completed

# === Gemini Setup ===
API_KEY = "AIzaSyAwaxJNzK4nviEVVFb8R8c_UGH0a9YbL1w"  # Replace with your Gemini API key
genai.configure(api_key=API_KEY)
model = genai.GenerativeModel("gemini-2.5-flash")

# === Utility: Clean OCR Text ===
def clean_ocr_text(text: Optional[str]) -> str:
    if not text:
        return ""
    text = re.sub(r'-\n', '', text)
    text = re.sub(r'\s*\n\s*', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'(\*\*|__)(.*?)\1', r'\2', text)
    text = re.sub(r'(\*|_)(.*?)\1', r'\2', text)
    return text.strip()

# === Core OCR Extraction Class ===
class FastPPTXOCRExtractor:
    def __init__(self, dpi: int = 300, max_workers: int = 4):
        self.dpi = dpi
        self.max_workers = max_workers

    def _extract_text_native(self, pptx_path: str) -> str:
        prs = Presentation(pptx_path)
        slides_text = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if slide_parts:
                slides_text.append(f"Slide {idx}:\n" + "\n".join(slide_parts))
        return "\n\n".join(slides_text)

    def _pptx_to_pdf(self, pptx_path: str, out_dir: str) -> str:
        subprocess.run([
            "soffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path
        ], check=True)
        base_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
        return os.path.join(out_dir, base_name)

    def _pdf_to_images(self, pdf_path: str, out_dir: str) -> list[str]:
        doc = fitz.open(pdf_path)
        paths = []
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(dpi=self.dpi)
            img_path = os.path.join(out_dir, f"slide_{i}.png")
            pix.save(img_path)
            paths.append(img_path)
        return paths

    def _gemini_ocr_image(self, image_path: str, prompt: str = "Extract all text from this slide:") -> str:
        with open(image_path, "rb") as f:
            image_data = f.read()
        response = model.generate_content([
            {"mime_type": "image/png", "data": image_data},
            prompt
        ])
        return clean_ocr_text(response.text)

    def _gemini_pdf_extract(self, pdf_path: str, prompt: str = "Extract slide-wise text from this presentation:") -> str:
        with open(pdf_path, "rb") as f:
            pdf_data = f.read()
        response = model.generate_content([
            {"mime_type": "application/pdf", "data": pdf_data},
            prompt
        ])
        return clean_ocr_text(response.text)

    def _parallel_ocr_images(self, image_paths: list[str]) -> list[str]:
        results = []
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            futures = [executor.submit(self._gemini_ocr_image, path) for path in image_paths]
            for future in as_completed(futures):
                try:
                    results.append(future.result())
                except Exception as e:
                    results.append(f"[ERROR] {e}")
        return results

    def extract_text_from_pptx_slides(self, pptx_path: str) -> list[str]:
        with tempfile.TemporaryDirectory() as tmp:
            native_text = self._extract_text_native(pptx_path)
            pdf_path = self._pptx_to_pdf(pptx_path, tmp)

            try:
                print("🔍 Using Gemini Vision Pro on full PDF...")
                gemini_text = self._gemini_pdf_extract(pdf_path)
                return [native_text, gemini_text]
            except Exception as e:
                print(f"⚠️ Gemini PDF extract failed: {e}")
                print("🔄 Falling back to image-based OCR...")
                images = self._pdf_to_images(pdf_path, tmp)
                ocr_texts = self._parallel_ocr_images(images)
                return [native_text] + ocr_texts
